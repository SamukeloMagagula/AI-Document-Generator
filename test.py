import os
import sys
import tempfile
import threading
import webbrowser
import json
from datetime import datetime
from http.server import BaseHTTPRequestHandler, HTTPServer
from urllib.parse import unquote_plus, urlparse, parse_qs

# ------------------------------
# NEW: LLM INTEGRATION (OpenAI + Grok compatible)
# ------------------------------
try:
    from openai import OpenAI
except ImportError:
    print("Please install: pip install openai")

# ==================== LLM CONFIGURATION ====================
# Best practice: set these as environment variables instead of hard-coding
USE_LLM = True
API_PROVIDER = "openai"          # Change to "grok" if you prefer xAI
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")   # or put your key here
GROK_API_KEY = os.getenv("XAI_API_KEY")        # or put your key here

def get_llm_client():
    if API_PROVIDER == "grok":
        return OpenAI(
            api_key=GROK_API_KEY or "dummy",
            base_url="https://api.x.ai/v1"
        )
    else:  # openai
        return OpenAI(api_key=OPENAI_API_KEY)

# ------------------------------
# LIBRARIES FOR REAL DOCUMENTS (unchanged)
# ------------------------------
try:
    from docx import Document
    from docx.shared import Pt, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
except ImportError:
    print("Please install: pip install python-docx")

try:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill
except ImportError:
    print("Please install: pip install openpyxl")

try:
    from pptx import Presentation
except ImportError:
    print("Please install: pip install python-pptx")

# ------------------------------
# LLM POWERED CONTENT GENERATION (the real magic)
# ------------------------------
def generate_llm_content(topic, doc_type="word", style="Professional"):
    topic = topic.strip().title()

    prompt = f"""You are an expert professional document creator.
Create a high-quality, detailed, {style.lower()} document about "{topic}" suitable for a {doc_type.upper()} file.

Return **ONLY** valid JSON with this exact structure (no extra text):
{{
  "title": "{topic}",
  "sections": [
    {{
      "title": "Section Title 1",
      "content": "Full paragraph here.",
      "bullets": ["bullet point 1", "bullet point 2", ...]
    }},
    ... (5 to 8 sections)
  ]
}}

Make the content accurate, engaging, and perfectly tailored to {style} style. Use real-world examples where appropriate."""

    try:
        client = get_llm_client()
        response = client.chat.completions.create(
            model="grok-4" if API_PROVIDER == "grok" else "gpt-4o-mini",
            messages=[
                {"role": "system", "content": "You are a precise document architect. Always respond with clean JSON only."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7,
            max_tokens=4000
        )

        content = response.choices[0].message.content.strip()
        # Extract JSON if model added extra text
        if content.startswith("```json"):
            content = content.split("```json")[1].split("```")[0]
        elif content.startswith("```"):
            content = content.split("```")[1]

        data = json.loads(content)
        return data["title"], data["sections"], style

    except Exception as e:
        print(f"LLM call failed: {e}. Falling back to template...")
        # Fallback to your original template
        return generate_template_content(topic, doc_type, style)


def generate_template_content(topic, doc_type="word", style="Professional"):
    # Your original sophisticated template (kept as safety net)
    sections = [
        {"title": "Introduction", "content": f"{topic} is a vital topic in today's world."},
        {"title": "Key Concepts", "content": f"Core principles and fundamentals of {topic}."},
        {"title": "Real-World Applications", "content": f"Practical uses across industries."},
        {"title": "Challenges", "content": "Common obstacles and solutions."},
        {"title": "Future Outlook", "content": f"Emerging trends in {topic}."},
        {"title": "Conclusion", "content": f"Summary of {topic}."}
    ]
    if style == "Creative":
        for s in sections:
            s["content"] += " Explored with fresh and inspiring perspectives."
    return topic, sections, style


# ------------------------------
# REAL DOCUMENT CREATORS (slightly improved for LLM output)
# ------------------------------
def create_real_word_doc(topic, sections, style):
    doc = Document()
    doc.add_heading(f"{topic} Report", 0)
    doc.add_paragraph(f"Generated on {datetime.now().strftime('%B %d, %Y')} • {style} Style • Powered by AI", style='Subtitle')

    for section in sections:
        doc.add_heading(section.get("title", "Section"), level=1)
        doc.add_paragraph(section.get("content", ""))
        for bullet in section.get("bullets", []):
            doc.add_paragraph(bullet, style='List Bullet')
        doc.add_paragraph()

    filepath = os.path.join(tempfile.gettempdir(), f"{topic.replace(' ', '_')}_Report.docx")
    doc.save(filepath)
    return filepath, os.path.basename(filepath)


def create_real_excel_doc(topic, sections, style):
    wb = Workbook()
    ws = wb.active
    ws.title = "Overview"
    ws['A1'] = f"{topic} Analysis"
    ws['A1'].font = Font(size=16, bold=True, color="FFFFFF")
    ws['A1'].fill = PatternFill(start_color="1e40af", end_color="1e40af", fill_type="solid")

    row = 3
    for section in sections:
        ws.cell(row=row, column=1, value=section.get("title")).font = Font(bold=True)
        row += 1
        ws.cell(row=row, column=1, value=section.get("content"))
        row += 1
        for bullet in section.get("bullets", []):
            ws.cell(row=row, column=2, value="• " + bullet)
            row += 1
        row += 1

    filename = f"{topic.replace(' ', '_')}_Analysis.xlsx"
    filepath = os.path.join(tempfile.gettempdir(), filename)
    wb.save(filepath)
    return filepath, filename


def create_real_ppt_doc(topic, sections, style):
    prs = Presentation()
    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = topic
    slide.placeholders[1].text = f"AI Generated • {style} Style"

    for section in sections:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = section.get("title", "")
        tf = slide.placeholders[1].text_frame
        tf.text = section.get("content", "")
        for bullet in section.get("bullets", []):
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0

    filename = f"{topic.replace(' ', '_')}_Presentation.pptx"
    filepath = os.path.join(tempfile.gettempdir(), filename)
    prs.save(filepath)
    return filepath, filename


# ------------------------------
# NATURAL LANGUAGE PARSER (unchanged – still excellent)
# ------------------------------
def parse_command(command):
    cmd_lower = command.lower().strip()
    if any(word in cmd_lower for word in ["powerpoint", "presentation", "slides", "ppt", "deck"]):
        doc_type = "ppt"
    elif any(word in cmd_lower for word in ["excel", "spreadsheet", "table", "sheet", "data", "budget", "forecast"]):
        doc_type = "excel"
    else:
        doc_type = "word"

    style = "Professional"
    for s in ["modern", "creative", "technical", "academic", "formal"]:
        if s in cmd_lower:
            style = s.capitalize()
            break

    # Extract topic
    triggers = ["about ", "on ", "regarding ", "for ", "create ", "make ", "generate "]
    topic = "General Topic"
    for trigger in triggers:
        if trigger in cmd_lower:
            topic = cmd_lower.split(trigger, 1)[1].split(" in ")[0].strip()
            break

    return doc_type, topic.title(), style


# ------------------------------
# HTTP HANDLER
# ------------------------------
class Handler(BaseHTTPRequestHandler):
    def do_GET(self):
        parsed = urlparse(self.path)

        if parsed.path == "/":
            self.send_response(200)
            self.send_header("Content-type", "text/html")
            self.end_headers()
            self.wfile.write(bytes(get_main_page(), "utf-8"))

        elif parsed.path.startswith("/generate"):
            query_params = parse_qs(parsed.query)
            cmd = unquote_plus(query_params.get("cmd", [""])[0]).strip()

            if not cmd:
                self.send_response(400)
                self.end_headers()
                self.wfile.write(b"Error: No command provided")
                return

            doc_type, topic, style = parse_command(cmd)

            # LLM magic happens here
            title, sections, tone = generate_llm_content(topic, doc_type, style)

            if doc_type == "excel":
                filepath, filename = create_real_excel_doc(title, sections, style)
                mime = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            elif doc_type == "ppt":
                filepath, filename = create_real_ppt_doc(title, sections, style)
                mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            else:
                filepath, filename = create_real_word_doc(title, sections, style)
                mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"

            self.send_response(200)
            self.send_header("Content-type", mime)
            self.send_header("Content-Disposition", f'attachment; filename="{filename}"')
            self.end_headers()

            with open(filepath, "rb") as f:
                self.wfile.write(f.read())

            try:
                os.remove(filepath)
            except:
                pass

        else:
            self.send_response(404)
            self.end_headers()
            self.wfile.write(b"404 - Not Found")


# ------------------------------
# MAIN UI – Now proudly "OfficeGPT style"
# ------------------------------
def get_main_page():
    return f"""
    <!DOCTYPE html>
    <html>
    <head>
        <meta charset="UTF-8">
        <title>OfficeGPT • Real AI Documents</title>
        <style>
            body {{ font-family: 'Segoe UI', system-ui, sans-serif; background: linear-gradient(135deg, #1e3a8a 0%, #3b82f6 100%); color: white; margin:0; padding:0; min-height:100vh; display:flex; align-items:center; justify-content:center; }}
            .container {{ max-width: 720px; width:90%; background:rgba(255,255,255,0.1); backdrop-filter:blur(16px); border-radius:24px; padding:60px 40px; box-shadow:0 20px 40px rgba(0,0,0,0.2); text-align:center; }}
            h1 {{ font-size:2.8em; margin-bottom:8px; }}
            p.tagline {{ font-size:1.3em; opacity:0.9; margin-bottom:40px; }}
            input {{ width:100%; padding:20px; font-size:1.25em; border:none; border-radius:16px; margin-bottom:20px; background:rgba(255,255,255,0.95); color:#1e2937; }}
            button {{ background:#fff; color:#1e40af; padding:18px 60px; font-size:1.3em; font-weight:bold; border:none; border-radius:16px; cursor:pointer; transition:all 0.2s; }}
            button:hover {{ transform:translateY(-3px); box-shadow:0 10px 20px rgba(0,0,0,0.2); }}
            .examples {{ margin-top:40px; font-size:0.95em; opacity:0.85; line-height:1.6; }}
        </style>
    </head>
    <body>
        <div class="container">
            <h1>📄 OfficeGPT</h1>
            <p class="tagline">Just say what you want — real LLM creates a beautiful document instantly</p>
            
            <input id="cmd" placeholder="Create a modern PowerPoint about the future of AI" autofocus />
            <br><br>
            <button onclick="generate()">Generate Document ✨</button>
            
            <div class="examples">
                <strong>Try these examples:</strong><br>
                • Make a presentation on climate change in creative style<br>
                • Excel budget forecast for a startup<br>
                • Professional Word report about renewable energy<br>
                • Technical slides explaining quantum computing
            </div>
        </div>

        <script>
        function generate() {{
            let cmd = document.getElementById("cmd").value.trim();
            if (!cmd) {{ alert("Just tell me what document you want!"); return; }}
            window.location.href = "/generate?cmd=" + encodeURIComponent(cmd);
        }}
        document.getElementById("cmd").addEventListener("keypress", e => {{ if (e.key === "Enter") generate(); }});
        </script>
    </body>
    </html>
    """


# ------------------------------
# START SERVER
# ------------------------------
def run_server():
    port = 8000
    server = HTTPServer(("localhost", port), Handler)
    print(f"🚀 OfficeGPT with REAL LLM running at http://localhost:{port}")
    print("✅ Just describe any document — GPT/Grok will create real content!")
    if not (OPENAI_API_KEY or GROK_API_KEY):
        print("⚠️  No API key detected → using template fallback (set OPENAI_API_KEY or XAI_API_KEY)")
    server.serve_forever()


if __name__ == "__main__":
    try:
        thread = threading.Thread(target=run_server, daemon=True)
        thread.start()
        webbrowser.open("http://localhost:8000")
        print("Press Ctrl+C to stop")
        while True:
            input()
    except KeyboardInterrupt:
        print("\nServer stopped.")
        sys.exit(0)